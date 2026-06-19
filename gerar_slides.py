from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Paleta de cores BrasilSeg / Corporativa
AZUL_ESCURO  = RGBColor(0x00, 0x2B, 0x5C)   # azul marinha
AZUL_MEDIO   = RGBColor(0x00, 0x56, 0xA2)   # azul corporativo
VERDE        = RGBColor(0x00, 0x9A, 0x44)   # verde BrasilSeg
AMARELO      = RGBColor(0xFF, 0xC2, 0x00)   # amarelo destaque
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO  = RGBColor(0xF2, 0xF2, 0xF2)
CINZA_TEXTO  = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def bg_fill(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False,
                color=BRANCO, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_multiline(slide, lines, l, t, w, h, font_size=16, color=BRANCO,
                  bold_first=False, bullet=False, line_spacing=1.2):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = "• " if bullet else ""
        run.text = prefix + line
        run.font.size = Pt(font_size)
        run.font.bold = (bold_first and i == 0)
        run.font.color.rgb = color
        run.font.name = "Calibri"

# ──────────────────────────────────────────────
# SLIDE 1 — CAPA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, AZUL_ESCURO)

# Faixa lateral verde
add_rect(slide, 0, 0, 0.5, 7.5, VERDE)

# Faixa amarela decorativa
add_rect(slide, 0.5, 5.8, 12.83, 0.08, AMARELO)

# Logo texto BrasilSeg
add_textbox(slide, "BRASILSEG", 1.0, 0.4, 5, 0.7, font_size=13,
            bold=True, color=AMARELO)

# Título principal
add_textbox(slide, "Modelo Funil Ágil\nde Produtos", 1.0, 1.3, 11, 2.2,
            font_size=46, bold=True, color=BRANCO)

# Subtítulo
add_textbox(slide, "Agilidade com Governança: Da Ideação à Aprovação\ncom Velocidade, Transparência e Segurança",
            1.0, 3.5, 11, 1.2, font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF))

# Linha separadora
add_rect(slide, 1.0, 4.9, 4, 0.05, AMARELO)

# Apresentadora
add_textbox(slide, "Andressa Marquito  |  Gerente de Portfolio de Projetos",
            1.0, 5.1, 9, 0.5, font_size=14, color=BRANCO)
add_textbox(slide, "Candidatura — Diretoria de Projetos e Agilidade",
            1.0, 5.6, 9, 0.4, font_size=12, color=RGBColor(0xAA, 0xBB, 0xDD))

# Ramos
add_textbox(slide, "Danos  |  Vida / Prestamista  |  Rural",
            1.0, 6.5, 8, 0.5, font_size=13, italic=True,
            color=RGBColor(0xAA, 0xCC, 0xFF))

# ──────────────────────────────────────────────
# SLIDE 2 — CENÁRIO ATUAL
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, CINZA_CLARO)
add_rect(slide, 0, 0, 0.5, 7.5, AZUL_ESCURO)
add_rect(slide, 0.5, 0, 12.83, 1.0, AZUL_ESCURO)

add_textbox(slide, "01  |  O CENÁRIO ATUAL", 1.0, 0.15, 10, 0.7,
            font_size=22, bold=True, color=BRANCO)

# 3 caixas de problema
boxes = [
    ("⏳ Velocidade", "Times ágeis aguardam filas de aprovação em diretoria,\nperdendo janelas de mercado."),
    ("⚖️ Priorização", "Competição subjetiva entre Danos, Vida/Prestamista\ne Rural por recursos e atenção executiva."),
    ("📋 Governança", "A avaliação de viabilidade econômica ocorre tarde,\naumentando retrabalho e desperdício."),
]
for i, (titulo, desc) in enumerate(boxes):
    x = 1.0 + i * 4.1
    add_rect(slide, x, 1.3, 3.8, 3.2, AZUL_ESCURO)
    add_textbox(slide, titulo, x + 0.2, 1.5, 3.4, 0.7,
                font_size=17, bold=True, color=AMARELO)
    add_textbox(slide, desc, x + 0.2, 2.3, 3.4, 2.0,
                font_size=14, color=BRANCO)

# Frase de impacto
add_rect(slide, 1.0, 4.8, 11.33, 1.4, AZUL_MEDIO)
add_textbox(slide,
            "Resultado: Metodologia ágil implantada, mas sem entrega de valor real.\n"
            "O problema não é a metodologia — é o modelo de operação.",
            1.3, 4.9, 10.8, 1.2, font_size=16, bold=True, color=BRANCO,
            align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 3 — A PROPOSTA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, AZUL_ESCURO)
add_rect(slide, 0, 0, 0.5, 7.5, VERDE)
add_rect(slide, 0.5, 0, 12.83, 1.0, RGBColor(0x00, 0x3A, 0x70))

add_textbox(slide, "02  |  A PROPOSTA", 1.0, 0.15, 10, 0.7,
            font_size=22, bold=True, color=BRANCO)

add_textbox(slide, "Modelo Funil Ágil de Produtos BrasilSeg",
            1.0, 1.2, 11.3, 0.8, font_size=30, bold=True, color=AMARELO)

add_textbox(slide,
            "Um modelo de três fases com portões de decisão progressivos,\n"
            "onde Governança e Agilidade trabalham juntas — não contra si mesmas.",
            1.0, 2.1, 11.3, 0.9, font_size=17, color=RGBColor(0xCC, 0xDD, 0xFF))

# 3 fases
fases = [
    ("01", "DESCOBERTA", "2 a 4 semanas", VERDE,
     ["Lean Canvas do produto", "Análise de mercado e concorrência",
      "Checklist regulatório SUSEP", "Validação com área técnica do ramo"]),
    ("02", "VALIDAÇÃO", "4 a 6 semanas", AZUL_MEDIO,
     ["MVP conceitual / protótipo", "Viabilidade econômica preliminar",
      "Análise de risco e compliance", "Aprovação no Comitê Técnico"]),
    ("03", "APROVAÇÃO & ENTREGA", "A partir da semana 10", RGBColor(0xC0, 0x39, 0x2B),
     ["Business case completo", "Apresentação à Diretoria", "Sprints de desenvolvimento",
      "Monitoramento de OKRs"]),
]
for i, (num, nome, prazo, cor, items) in enumerate(fases):
    x = 1.0 + i * 4.1
    add_rect(slide, x, 3.1, 3.8, 3.8, cor)
    add_textbox(slide, num, x + 0.15, 3.15, 0.6, 0.55,
                font_size=28, bold=True, color=BRANCO)
    add_textbox(slide, nome, x + 0.15, 3.7, 3.5, 0.5,
                font_size=13, bold=True, color=BRANCO)
    add_textbox(slide, prazo, x + 0.15, 4.1, 3.5, 0.35,
                font_size=11, italic=True, color=RGBColor(0xFF, 0xFF, 0xAA))
    add_multiline(slide, items, x + 0.15, 4.5, 3.5, 2.2,
                  font_size=11, color=BRANCO, bullet=True)

# ──────────────────────────────────────────────
# SLIDE 4 — PRIORIZAÇÃO DE PORTFÓLIO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, CINZA_CLARO)
add_rect(slide, 0, 0, 0.5, 7.5, AZUL_ESCURO)
add_rect(slide, 0.5, 0, 12.83, 1.0, AZUL_ESCURO)

add_textbox(slide, "03  |  PRIORIZAÇÃO DE PORTFÓLIO", 1.0, 0.15, 11, 0.7,
            font_size=22, bold=True, color=BRANCO)

add_textbox(slide, "Uma fila única, objetiva e auditável para os 3 Ramos",
            1.0, 1.1, 11, 0.55, font_size=19, bold=True, color=AZUL_ESCURO)

# Tabela de critérios
criterios = [
    ("Critério", "Peso", "Como é medido"),
    ("Potencial de Receita", "30%", "Projeção de prêmio emitido no primeiro ano"),
    ("Alinhamento Estratégico", "25%", "Aderência ao plano estratégico BrasilSeg"),
    ("Complexidade Regulatória", "20%", "Classificação SUSEP: baixa / média / alta"),
    ("Esforço de Implementação", "25%", "Estimativa de squads e prazo de entrega"),
]
header_colors = [AZUL_ESCURO, AZUL_ESCURO, AZUL_ESCURO]
row_colors    = [CINZA_CLARO, BRANCO]

for r, row in enumerate(criterios):
    y = 1.9 + r * 0.75
    bg = AZUL_ESCURO if r == 0 else (CINZA_CLARO if r % 2 == 0 else BRANCO)
    add_rect(slide, 1.0, y, 10.8, 0.72, bg)
    widths = [4.5, 1.8, 4.5]
    xs     = [1.0, 5.5, 7.3]
    fc = BRANCO if r == 0 else CINZA_TEXTO
    for c, (col, cw, cx) in enumerate(zip(row, widths, xs)):
        add_textbox(slide, col, cx + 0.1, y + 0.15, cw - 0.2, 0.5,
                    font_size=13 if r > 0 else 13,
                    bold=(r == 0), color=fc)

# Nota
add_textbox(slide,
            "★  Cada ramo submete seus produtos com pontuação preenchida."
            "  A fila é ordenada automaticamente — decisão transparente para toda a diretoria.",
            1.0, 6.6, 11, 0.6, font_size=12, italic=True, color=AZUL_MEDIO)

# ──────────────────────────────────────────────
# SLIDE 5 — GESTÃO DE MUDANÇA
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, AZUL_ESCURO)
add_rect(slide, 0, 0, 0.5, 7.5, AMARELO)
add_rect(slide, 0.5, 0, 12.83, 1.0, RGBColor(0x00, 0x3A, 0x70))

add_textbox(slide, "04  |  GESTÃO DA MUDANÇA", 1.0, 0.15, 11, 0.7,
            font_size=22, bold=True, color=BRANCO)

add_textbox(slide, "Como levaremos as equipes e a liderança nessa jornada",
            1.0, 1.1, 11, 0.5, font_size=17, color=RGBColor(0xCC, 0xDD, 0xFF))

pilares = [
    ("🎯 Engajamento da Liderança",
     ["Workshops com diretores dos 3 ramos", "Patrocínio visível da presidência",
      "Comitê de Agilidade com representantes executivos"]),
    ("📚 Capacitação das Equipes",
     ["Treinamento em práticas ágeis (Scrum, Kanban)", "Formação de Agile Coaches internos",
      "Comunidade de prática entre ramos"]),
    ("📊 Comunicação e Transparência",
     ["Dashboard de portfólio visível para toda liderança", "Rituais mensais de revisão de portfólio",
      "Relatórios RAG por produto e ramo"]),
    ("🔄 Melhoria Contínua",
     ["Retrospectivas trimestrais do modelo", "Benchmarking com outras seguradoras",
      "Ajuste dos pesos da matriz de priorização"]),
]
for i, (titulo, items) in enumerate(pilares):
    col = i % 2
    row = i // 2
    x = 1.0 + col * 6.0
    y = 1.9 + row * 2.5
    add_rect(slide, x, y, 5.6, 2.2, RGBColor(0x00, 0x3A, 0x70))
    add_rect(slide, x, y, 5.6, 0.5, AZUL_MEDIO)
    add_textbox(slide, titulo, x + 0.15, y + 0.08, 5.3, 0.4,
                font_size=13, bold=True, color=BRANCO)
    add_multiline(slide, items, x + 0.15, y + 0.6, 5.3, 1.5,
                  font_size=12, color=RGBColor(0xCC, 0xDD, 0xFF), bullet=True)

# ──────────────────────────────────────────────
# SLIDE 6 — ROADMAP
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, CINZA_CLARO)
add_rect(slide, 0, 0, 0.5, 7.5, AZUL_ESCURO)
add_rect(slide, 0.5, 0, 12.83, 1.0, AZUL_ESCURO)

add_textbox(slide, "05  |  ROADMAP DE IMPLEMENTAÇÃO", 1.0, 0.15, 11, 0.7,
            font_size=22, bold=True, color=BRANCO)

etapas = [
    ("M1–M2", "Diagnóstico &\nCapacitação",
     AZUL_ESCURO, ["Mapeamento AS-IS", "Treinamento de líderes", "Formação dos squads"]),
    ("M3–M4", "Projeto Piloto",
     VERDE, ["1 produto piloto por ramo", "Aplicação do funil ágil", "Coleta de métricas"]),
    ("M5–M6", "Ajustes &\nEscala",
     AZUL_MEDIO, ["Retrospectiva do piloto", "Refinamento do modelo", "Expansão para todos os ramos"]),
    ("M7+", "Operação\nPlena",
     RGBColor(0xC0, 0x39, 0x2B), ["Dashboard de portfólio ativo", "Governança ágil consolidada", "Melhoria contínua"]),
]
for i, (periodo, titulo, cor, items) in enumerate(etapas):
    x = 1.0 + i * 3.1
    add_rect(slide, x, 1.2, 2.9, 0.55, cor)
    add_textbox(slide, periodo, x + 0.1, 1.25, 2.7, 0.45,
                font_size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 1.75, 2.9, 1.3, RGBColor(0xE0, 0xE8, 0xF4))
    add_textbox(slide, titulo, x + 0.1, 1.8, 2.7, 1.2,
                font_size=14, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 3.05, 2.9, 2.5, BRANCO)
    add_multiline(slide, items, x + 0.15, 3.1, 2.6, 2.4,
                  font_size=12, color=CINZA_TEXTO, bullet=True)
    # seta
    if i < 3:
        add_textbox(slide, "▶", x + 3.0, 2.2, 0.3, 0.4,
                    font_size=16, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 7 — MÉTRICAS DE SUCESSO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, AZUL_ESCURO)
add_rect(slide, 0, 0, 0.5, 7.5, VERDE)
add_rect(slide, 0.5, 0, 12.83, 1.0, RGBColor(0x00, 0x3A, 0x70))

add_textbox(slide, "06  |  MÉTRICAS DE SUCESSO", 1.0, 0.15, 11, 0.7,
            font_size=22, bold=True, color=BRANCO)

add_textbox(slide, "O que vamos medir para garantir que o modelo funciona",
            1.0, 1.1, 11, 0.5, font_size=17, color=RGBColor(0xCC, 0xDD, 0xFF))

metricas = [
    ("⏱️", "Time-to-Approval", "Tempo médio de\nideação → diretoria",
     "6 meses", "10 semanas", VERDE),
    ("✅", "Taxa de Aprovação", "% de produtos aprovados\nna 1ª apresentação",
     "~50%", "> 80%", AZUL_MEDIO),
    ("🚀", "Produtos Lançados", "Novos produtos lançados\npor ano (3 ramos)",
     "Baseline atual", "+40%", RGBColor(0xC0, 0x39, 0x2B)),
    ("😊", "NPS Interno", "Satisfação das equipes\ncom o processo",
     "Não medido", "> 70", RGBColor(0xE6, 0x7E, 0x22)),
]
for i, (icon, nome, desc, antes, depois, cor) in enumerate(metricas):
    x = 1.0 + i * 3.1
    add_rect(slide, x, 1.8, 2.9, 4.5, RGBColor(0x00, 0x3A, 0x70))
    add_textbox(slide, icon, x + 1.2, 1.9, 0.6, 0.55,
                font_size=24, align=PP_ALIGN.CENTER)
    add_textbox(slide, nome, x + 0.1, 2.5, 2.7, 0.55,
                font_size=13, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x + 0.1, 3.05, 2.7, 0.75,
                font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.2, 3.85, 2.5, 0.6, RGBColor(0x60, 0x20, 0x20))
    add_textbox(slide, f"Hoje: {antes}", x + 0.2, 3.9, 2.5, 0.5,
                font_size=12, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.2, 4.55, 2.5, 0.6, cor)
    add_textbox(slide, f"Meta: {depois}", x + 0.2, 4.6, 2.5, 0.5,
                font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 8 — PRÓXIMOS PASSOS / CALL TO ACTION
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, CINZA_CLARO)
add_rect(slide, 0, 0, 0.5, 7.5, AZUL_ESCURO)
add_rect(slide, 0.5, 0, 12.83, 1.0, AZUL_ESCURO)

add_textbox(slide, "07  |  PRÓXIMOS PASSOS", 1.0, 0.15, 11, 0.7,
            font_size=22, bold=True, color=BRANCO)

add_textbox(slide, "O que precisamos para começar ainda neste semestre",
            1.0, 1.1, 11, 0.5, font_size=17, bold=False, color=AZUL_ESCURO)

passos = [
    ("1", "Aprovação da iniciativa", "Patrocínio da Presidência e endosso dos Diretores de Ramo",
     AZUL_ESCURO),
    ("2", "Formação do Comitê de Agilidade", "1 representante executivo por ramo + GP de Portfólio",
     VERDE),
    ("3", "Seleção do Produto Piloto", "1 produto por ramo para validação do modelo em condições reais",
     AZUL_MEDIO),
    ("4", "Kickoff do Projeto", "Workshop de 2 dias com líderes e squads para alinhamento e início",
     RGBColor(0xE6, 0x7E, 0x22)),
]
for i, (num, titulo, desc, cor) in enumerate(passos):
    y = 1.8 + i * 1.2
    add_rect(slide, 1.0, y, 0.65, 0.85, cor)
    add_textbox(slide, num, 1.0, y + 0.1, 0.65, 0.65,
                font_size=24, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.65, y, 10.68, 0.85, BRANCO)
    add_textbox(slide, titulo, 1.8, y + 0.05, 5, 0.4,
                font_size=14, bold=True, color=AZUL_ESCURO)
    add_textbox(slide, desc, 1.8, y + 0.45, 10, 0.35,
                font_size=12, color=CINZA_TEXTO)

add_rect(slide, 1.0, 6.6, 11.33, 0.6, AZUL_ESCURO)
add_textbox(slide,
            "\"Não se trata de ser mais ágil. Trata-se de ser mais inteligente na forma de decidir e entregar.\"",
            1.2, 6.63, 10.9, 0.55, font_size=13, italic=True,
            color=AMARELO, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 9 — ENCERRAMENTO
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_fill(slide, AZUL_ESCURO)
add_rect(slide, 0, 0, 0.5, 7.5, VERDE)
add_rect(slide, 0.5, 5.5, 12.83, 2.0, RGBColor(0x00, 0x1A, 0x40))

add_textbox(slide, "Obrigada.", 1.0, 1.5, 11, 1.5,
            font_size=60, bold=True, color=BRANCO)

add_textbox(slide,
            "Estou pronta para liderar essa transformação\n"
            "e construir, junto com a BrasilSeg, um modelo de portfólio\n"
            "que entrega mais produtos, com mais velocidade e mais confiança.",
            1.0, 3.0, 11, 2.0, font_size=18, color=RGBColor(0xCC, 0xDD, 0xFF))

add_rect(slide, 1.0, 5.2, 4.5, 0.05, AMARELO)

add_textbox(slide, "Andressa Marquito", 1.0, 5.5, 8, 0.5,
            font_size=16, bold=True, color=BRANCO)
add_textbox(slide, "dessajf@gmail.com", 1.0, 6.0, 8, 0.4,
            font_size=13, color=RGBColor(0xAA, 0xBB, 0xDD))

# ──────────────────────────────────────────────
# SALVAR
# ──────────────────────────────────────────────
output = r"c:\Users\dessa\OneDrive\Área de Trabalho\Funil_Agil_BrasilSeg.pptx"
prs.save(output)
print(f"Arquivo salvo: {output}")
